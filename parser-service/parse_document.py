"""
MinerU 文档解析服务（Flask REST API）
内置于智能情报分析平台，提供高质量文档解析能力

支持格式：
- PDF → Markdown（MinerU pipeline / marker）
- HTML → Markdown（MinerU-HTML / BeautifulSoup fallback）
- DOCX/PPTX → Markdown（MinerU）
- 图片 → Markdown（MinerU OCR）

启动方式：
  python parse_document.py --port 8100

Java 后端通过 HTTP 调用此服务：
  POST /api/parse  - 解析本地文件
  POST /api/parse-url - 解析URL网页
  GET /api/health - 健康检查
"""

import argparse
import json
import os
import sys
import tempfile
import traceback
from pathlib import Path

from flask import Flask, request, jsonify
from flask_cors import CORS

app = Flask(__name__)
CORS(app)

# 检测可用的解析器
import shutil as _shutil
_mineru_available = _shutil.which('mineru') is not None
_marker_available = False

if _mineru_available:
    print("[MinerU] mineru CLI found", file=sys.stderr)
else:
    # 检查 Python module
    try:
        import mineru
        _mineru_available = True
        print("[MinerU] mineru Python module available", file=sys.stderr)
    except ImportError:
        print("[MinerU] mineru NOT available", file=sys.stderr)

try:
    from marker.converters.pdf import PdfConverter
    _marker_available = True
    print("[MinerU] marker loaded successfully", file=sys.stderr)
except ImportError as e:
    print(f"[MinerU] marker not available: {e}", file=sys.stderr)

# HTML 解析 fallback
try:
    from bs4 import BeautifulSoup
    import re
    _bs4_available = True
except ImportError:
    _bs4_available = False


def parse_pdf_mineru(file_path: str, **kwargs) -> dict:
    """使用 MinerU CLI 解析 PDF/DOCX/PPTX 文件"""
    import subprocess
    import shutil

    # 创建临时输出目录
    output_dir = tempfile.mkdtemp(prefix='mineru_out_')

    try:
        # 使用 MinerU CLI：mineru -p <input> -o <output> -b pipeline
        mineru_bin = shutil.which('mineru')
        if not mineru_bin:
            raise RuntimeError("mineru command not found in PATH")

        cmd = [
            mineru_bin,
            '-p', file_path,
            '-o', output_dir,
            '-b', 'pipeline',  # pipeline backend 不需要 GPU
        ]

        print(f"[MinerU CLI] Running: {' '.join(cmd)}", file=sys.stderr)
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=600,  # 10分钟超时
        )

        if result.returncode != 0:
            print(f"[MinerU CLI] returncode={result.returncode}", file=sys.stderr)
            print(f"[MinerU CLI] stderr: {result.stderr[:1000]}", file=sys.stderr)
            print(f"[MinerU CLI] stdout: {result.stdout[:500]}", file=sys.stderr)

        # 读取输出的 markdown 文件
        markdown = ""
        images = []

        output_path = Path(output_dir)
        # MinerU 在 output_dir 下创建以文件名命名的子目录
        for md_file in output_path.rglob('*.md'):
            markdown += md_file.read_text(encoding='utf-8') + '\n'

        # 收集图片信息
        for img_file in output_path.rglob('*.png'):
            rel_path = str(img_file.relative_to(output_path))
            images.append({'path': rel_path, 'caption': ''})
        for img_file in output_path.rglob('*.jpg'):
            rel_path = str(img_file.relative_to(output_path))
            images.append({'path': rel_path, 'caption': ''})

        if not markdown:
            # 尝试读取 content_list.json
            for json_file in output_path.rglob('content_list.json'):
                with open(json_file) as f:
                    content_list = json.load(f)
                    for item in content_list:
                        if item.get('type') == 'text':
                            markdown += item.get('text', '') + '\n'
                        elif item.get('type') == 'table':
                            markdown += item.get('html', '') + '\n'
                        elif item.get('type') == 'image':
                            images.append({
                                'path': item.get('img_path', ''),
                                'caption': item.get('img_caption', ''),
                            })

        return {
            'markdown': markdown,
            'images': images,
            'parser': 'mineru-cli',
        }
    finally:
        # 清理临时目录
        shutil.rmtree(output_dir, ignore_errors=True)


def parse_pdf_marker(file_path: str) -> dict:
    """使用 marker 解析 PDF 文件（备选方案）"""
    if not _marker_available:
        raise RuntimeError("marker not installed. Run: pip install marker-pdf")

    converter = PdfConverter()
    result = converter(file_path)

    return {
        'markdown': result.markdown,
        'images': [],
        'parser': 'marker',
    }


def parse_html_content(html_content: str) -> dict:
    """解析 HTML 内容为 Markdown"""
    if not _bs4_available:
        # 简单 fallback：去标签
        import re
        text = re.sub(r'<[^>]+>', ' ', html_content)
        text = re.sub(r'\s+', ' ', text).strip()
        return {'markdown': text, 'parser': 'fallback'}

    soup = BeautifulSoup(html_content, 'html.parser')

    # 移除 script/style
    for tag in soup(['script', 'style', 'nav', 'footer', 'header']):
        tag.decompose()

    # 尝试提取主内容（article 或 main）
    main_content = soup.find('article') or soup.find('main') or soup.find('body') or soup
    if main_content is None:
        main_content = soup

    # 提取标题
    title = ''
    title_tag = soup.find('title')
    if title_tag:
        title = title_tag.get_text(strip=True)

    # 转换为 Markdown
    lines = []

    for element in main_content.descendants:
        if not hasattr(element, 'name'):
            continue

        if element.name == 'h1':
            lines.append(f"# {element.get_text(strip=True)}\n")
        elif element.name == 'h2':
            lines.append(f"## {element.get_text(strip=True)}\n")
        elif element.name == 'h3':
            lines.append(f"### {element.get_text(strip=True)}\n")
        elif element.name in ('h4', 'h5', 'h6'):
            level = int(element.name[1])
            lines.append(f"{'#' * level} {element.get_text(strip=True)}\n")
        elif element.name == 'p':
            text = element.get_text(strip=True)
            if text:
                lines.append(text + '\n')
        elif element.name == 'li':
            text = element.get_text(strip=True)
            if text:
                lines.append(f"- {text}\n")
        elif element.name == 'blockquote':
            text = element.get_text(strip=True)
            if text:
                lines.append(f"> {text}\n")
        elif element.name == 'code':
            text = element.get_text(strip=True)
            if text:
                lines.append(f"`{text}`\n")
        elif element.name == 'pre':
            text = element.get_text(strip=True)
            if text:
                lines.append(f"```\n{text}\n```\n")
        elif element.name == 'a':
            text = element.get_text(strip=True)
            href = element.get('href', '')
            if text and href:
                lines.append(f"[{text}]({href})\n")
        elif element.name == 'img':
            src = element.get('src', '')
            alt = element.get('alt', '')
            if src:
                lines.append(f"![{alt}]({src})\n")
        elif element.name == 'table':
            # 简单表格转 Markdown
            rows = element.find_all('tr')
            if rows:
                for i, row in enumerate(rows):
                    cells = row.find_all(['td', 'th'])
                    cell_texts = [c.get_text(strip=True) for c in cells]
                    lines.append('| ' + ' | '.join(cell_texts) + ' |')
                    if i == 0:
                        lines.append('| ' + ' | '.join(['---'] * len(cell_texts)) + ' |')
                lines.append('')
        elif element.name == 'br':
            lines.append('\n')

    # 去重和清理
    markdown = '\n'.join(lines)
    # 清理多余空行
    markdown = re.sub(r'\n{3,}', '\n\n', markdown)
    # 清理行内多余空格
    markdown = re.sub(r'  +', ' ', markdown)

    return {
        'markdown': markdown.strip(),
        'title': title,
        'parser': 'beautifulsoup',
    }


def parse_html_mineru(html_path: str) -> dict:
    """尝试使用 MinerU 解析 HTML"""
    try:
        converter = DocumentConverter()
        result = converter(html_path)
        markdown = ""
        if hasattr(result, 'md_content'):
            markdown = result.md_content
        elif hasattr(result, 'content_list'):
            for item in result.content_list:
                if item.get('type') == 'text':
                    markdown += item.get('text', '') + '\n'
        return {
            'markdown': markdown,
            'parser': 'mineru-html',
        }
    except Exception:
        # fallback to BeautifulSoup
        with open(html_path, 'r', encoding='utf-8', errors='ignore') as f:
            return parse_html_content(f.read())


@app.route('/api/health', methods=['GET'])
def health():
    return jsonify({
        'status': 'ok',
        'mineru': _mineru_available,
        'marker': _marker_available,
        'bs4': _bs4_available,
    })


@app.route('/api/parse', methods=['POST'])
def parse_file():
    """
    解析本地文件
    
    请求参数：
    - file_path: 文件路径
    - parser: 解析器选择 (mineru/marker/auto)，默认 auto
    
    返回：
    - markdown: 解析后的 Markdown 内容
    - images: 提取的图片列表
    - parser: 使用的解析器
    - title: 文档标题（HTML 文档）
    """
    data = request.get_json()
    if not data or 'file_path' not in data:
        return jsonify({'error': 'file_path is required'}), 400

    file_path = data['file_path']
    parser_choice = data.get('parser', 'auto')

    if not os.path.exists(file_path):
        return jsonify({'error': f'File not found: {file_path}'}), 404

    ext = Path(file_path).suffix.lower()

    try:
        if ext == '.pdf':
            if parser_choice == 'marker' and _marker_available:
                result = parse_pdf_marker(file_path)
            elif _mineru_available:
                result = parse_pdf_mineru(file_path)
            elif _marker_available:
                result = parse_pdf_marker(file_path)
            else:
                # 最终 fallback：尝试 PyPDF2 或 pdfminer
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(file_path)
                    text = '\n'.join([page.extract_text() or '' for page in reader.pages])
                    result = {'markdown': text, 'parser': 'pypdf'}
                except ImportError:
                    raise RuntimeError("No PDF parser available. Install mineru or marker-pdf or pypdf")

        elif ext in ('.html', '.htm'):
            if parser_choice == 'mineru' and _mineru_available:
                result = parse_html_mineru(file_path)
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    result = parse_html_content(f.read())

        elif ext in ('.docx', '.doc'):
            if _mineru_available:
                result = parse_pdf_mineru(file_path)
            else:
                try:
                    from docx import Document as DocxDocument
                    doc = DocxDocument(file_path)
                    text = '\n'.join([p.text for p in doc.paragraphs])
                    result = {'markdown': text, 'parser': 'python-docx'}
                except ImportError:
                    raise RuntimeError("No DOCX parser available. Install mineru or python-docx")

        elif ext in ('.pptx', '.ppt'):
            if _mineru_available:
                result = parse_pdf_mineru(file_path)
            else:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text_parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'):
                                text_parts.append(shape.text)
                    result = {'markdown': '\n'.join(text_parts), 'parser': 'python-pptx'}
                except ImportError:
                    raise RuntimeError("No PPTX parser available. Install mineru or python-pptx")

        elif ext in ('.xlsx', '.xls'):
            if _mineru_available:
                result = parse_pdf_mineru(file_path)
            else:
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path)
                    md_parts = []
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        md_parts.append(f"## {sheet}\n")
                        for row in ws.iter_rows(values_only=True):
                            cells = [str(c) if c is not None else '' for c in row]
                            md_parts.append('| ' + ' | '.join(cells) + ' |')
                        md_parts.append('')
                    result = {'markdown': '\n'.join(md_parts), 'parser': 'openpyxl'}
                except ImportError:
                    raise RuntimeError("No XLSX parser available. Install mineru or openpyxl")

        elif ext in ('.txt', '.md', '.markdown', '.csv'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                result = {'markdown': f.read(), 'parser': 'plaintext'}

        else:
            return jsonify({'error': f'Unsupported file type: {ext}'}), 400

        return jsonify(result)

    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': str(e), 'traceback': traceback.format_exc()}), 500


@app.route('/api/parse-url', methods=['POST'])
def parse_url():
    """
    解析URL网页内容
    
    请求参数：
    - url: 网页URL
    
    返回：
    - markdown: 解析后的 Markdown 内容
    - title: 网页标题
    - parser: 使用的解析器
    """
    data = request.get_json()
    if not data or 'url' not in data:
        return jsonify({'error': 'url is required'}), 400

    url = data['url']

    try:
        import urllib.request
        req = urllib.request.Request(
            url,
            headers={'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36'}
        )
        with urllib.request.urlopen(req, timeout=30) as response:
            html_content = response.read().decode('utf-8', errors='ignore')

        result = parse_html_content(html_content)
        result['url'] = url
        return jsonify(result)

    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/api/parse-html', methods=['POST'])
def parse_html_direct():
    """
    直接解析 HTML 内容（不需要文件）
    
    请求参数：
    - html: HTML 字符串
    
    返回：
    - markdown: 解析后的 Markdown 内容
    - title: 网页标题
    """
    data = request.get_json()
    if not data or 'html' not in data:
        return jsonify({'error': 'html content is required'}), 400

    try:
        result = parse_html_content(data['html'])
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='MinerU Document Parser Service')
    parser.add_argument('--port', type=int, default=8100, help='Port to listen on')
    parser.add_argument('--host', type=str, default='127.0.0.1', help='Host to bind to')
    args = parser.parse_args()

    print(f"[MinerU Service] Starting on {args.host}:{args.port}")
    print(f"  MinerU: {'available' if _mineru_available else 'NOT available'}")
    print(f"  marker: {'available' if _marker_available else 'NOT available'}")
    print(f"  bs4: {'available' if _bs4_available else 'NOT available'}")

    app.run(host=args.host, port=args.port, debug=False)
